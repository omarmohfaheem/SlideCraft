import os
# Module imports
import text_summarization as ts
import header_generation as hg
import image_generation as ig
import image_searching as ss
import read_file as rf

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_ALIGN

# Image imports
from io import BytesIO
import requests
from PIL import Image, UnidentifiedImageError

def remove_newlines(text):
    return text.replace('\n', ' ')

def read_document(file_path):
    """Read the content of the document."""
    try:
        return rf.read_file(file_path)
    except Exception as e:
        print(f"Error reading file: {e}")
        return None

def read_input(input_data, is_file=True):
    """Read the input data, either from a file or directly from text."""
    if is_file:
        return read_document(input_data)
    else:
        return input_data

def summarize_content(content):
    """Summarize the document content into bullet points."""
    content = remove_newlines(content)
    return ts.summarize(content)

def calculate_chunk_size(total_bullet_points, num_slides, min_bullet_points_per_slide=3):
    """Calculate the chunk size for distributing bullet points across slides."""
    if total_bullet_points <= num_slides * min_bullet_points_per_slide:
        return min_bullet_points_per_slide
    else:
        return (total_bullet_points + num_slides - 1) // num_slides  # Evenly distribute bullet points

def create_presentation(theme):
    """Create a new PowerPoint presentation with the specified theme."""
    return Presentation(f"../themes/{theme}.pptx")

def add_new_slide(prs, header ,theme, image_stream=None):
    """Add a new slide to the presentation with the given header and image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_shape = slide.shapes.title
    title_shape.text = header

     # Set title color based on the theme
    if theme == "Dark":
        title_color = RGBColor(255, 255, 255)
        text_align = PP_ALIGN.CENTER
    elif theme == "Creative":
        title_color = RGBColor(0, 112, 192)
        text_align = PP_ALIGN.CENTER
    elif theme == "Modern":
        title_color = RGBColor(0, 0, 0)
        text_align = PP_ALIGN.CENTER
    elif theme == "Abstract":
        title_color = RGBColor(112, 48, 160)
        text_align = PP_ALIGN.CENTER
    elif theme == "Technology":
        title_color = RGBColor(255, 0, 102)
        text_align = PP_ALIGN.RIGHT
    elif theme == "Business":
        title_color = RGBColor(16, 37, 63)
        text_align = PP_ALIGN.RIGHT
    else:
        title_color = RGBColor(0, 0, 0)  # Default color

    title_shape.text_frame.paragraphs[0].alignment = text_align # Align text to the right
    title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = title_color
    title_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(40)

    # Add a text box for bullet points
    left = Inches(0.2)
    top = Inches(1)
    width = Inches(6.7)
    height = Inches(4.0)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    # Add the generated image to the slide
    if image_stream:
        img_left = Inches(7)
        img_top = Inches(2)
        img_height = Inches(2.5)
        img_width = Inches(2.5)
        slide.shapes.add_picture(image_stream, img_left, img_top, height=img_height, width=img_width)

    return slide, text_frame

def add_dot_bullet(paragraph):
    """Add a dot bullet to the paragraph."""
    pPr = paragraph._p.get_or_add_pPr()
    buChar = OxmlElement('a:buChar')
    buChar.set('char', '•')
    pPr.insert(0, buChar)

def add_bullet_points(text_frame, chunk):
    """Add bullet points to the text frame."""
    for point in chunk:
        if point.strip():
            p = text_frame.add_paragraph()
            p.text = ' ' + point  
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 0, 0)
            add_dot_bullet(p)
            p.space_after = Pt(14)  
            p.space_before = Pt(14)
    
def generate_image(bullet_points_chunk, image_source):
    """Generate an image based on the bullet points chunk using the specified source."""
    header = " ".join(bullet_points_chunk)
    
    if image_source == "openai":
        prompt = f"Generate an image based on the following bullet points: {header}"
        image_url = ig.generate_image_from_text(prompt)
        image_response = requests.get(image_url)
        image = Image.open(BytesIO(image_response.content))

    elif image_source == "google":
        search_query = header
        image_urls = ss.search_google_images(search_query)
        for image_url in image_urls:
            print(f"Trying image URL: {image_url}")
            try:
                image_response = requests.get(image_url)
                image_response.raise_for_status()
                image = Image.open(BytesIO(image_response.content))
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                break
            except (requests.RequestException, UnidentifiedImageError) as e:
                print(f"Error generating image from URL {image_url}: {e}")
                continue

    image_stream = BytesIO()
    image.save(image_stream, format='JPEG')
    image_stream.seek(0)
    return image_stream

def generate_slides(title, input_data, is_file, num_slides, additional_slides, use_images, image_source, theme):
    """Generate a PowerPoint presentation from the given document or text."""
    text_content = read_input(input_data, is_file)
    if text_content is None:
        return None

    bullet_points = summarize_content(text_content)
    total_bullet_points = len(bullet_points)
    chunk_size = calculate_chunk_size(total_bullet_points, num_slides)

    prs = create_presentation(theme)

    # Add title slide if requested
    if additional_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = title
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"This Presentation is about {title}"

    # Process bullet points and generate slides
    bullet_points_index = 0
    for i in range(num_slides):
        if bullet_points_index < total_bullet_points:
            chunk = bullet_points[bullet_points_index:bullet_points_index + chunk_size]
            bullet_points_index += chunk_size
        else:
            chunk = []

        header = hg.generate_header(chunk) if chunk else "Slide Title"

        # Generate an image for the current chunk
        image_stream = None
        if use_images and chunk:
            try:
                image_stream = generate_image(chunk, image_source)
            except Exception as e:
                print(f"Error generating image: {e}")

        # Add a new slide for each chunk with the generated image
        slide, text_frame = add_new_slide(prs, header, theme, image_stream)

        # Add bullet points with smaller font size and dot bullet
        add_bullet_points(text_frame, chunk)

    # Add empty slides if there are more slides than bullet points
    while i < num_slides - 1:
        i += 1
        header = "Slide Title"
        add_new_slide(prs, header, theme)


    # Add thank you slide if requested
    if additional_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = "Thank You"
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = "Any Questions?"

    # Save the presentation to a BytesIO object
    presentation_stream = BytesIO()
    prs.save(presentation_stream)
    presentation_stream.seek(0)
    print("Presentation created successfully!")
    
    return presentation_stream

